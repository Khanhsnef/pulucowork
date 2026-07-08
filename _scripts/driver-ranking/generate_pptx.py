import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. Khởi tạo Presentation với tỷ lệ 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Bảng màu chuẩn Ahamove & Corporate
NAVY = RGBColor(14, 65, 116)      # #0E4174 (Ahamove Navy)
NAVY_LIGHT = RGBColor(26, 90, 154) # #1A5A9A
ORANGE = RGBColor(255, 127, 50)   # #FF7F32 (Ahamove Orange)
ORANGE_BG = RGBColor(255, 247, 237) # #FFF7ED
GREEN = RGBColor(16, 185, 129)    # #10B981
GREEN_BG = RGBColor(209, 250, 229) # #D1FAE5
RED = RGBColor(239, 68, 68)       # #EF4444
GRAY_TEXT = RGBColor(75, 85, 99)   # #4B5563
DARK_TEXT = RGBColor(17, 24, 39)   # #111827
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(249, 250, 251) # #F9FAFB
BORDER_COLOR = RGBColor(229, 231, 235) # #E5E7EB

def add_header(slide, title_text, category_text="DRIVERS LIFE-CYCLE STRATEGY"):
    """Thêm tiêu đề trang chuẩn hóa"""
    # Category Label
    cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = 'Arial'
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ORANGE

    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(12), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = NAVY

    # Orange line decoration
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.35), Inches(1.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.color.rgb = ORANGE

# ==========================================
# SLIDE 1: TIÊU ĐỀ DECK (COVER SLIDE) - NAVY THEME
# ==========================================
slide_layout = prs.slide_layouts[6] # Blank
slide = prs.slides.add_slide(slide_layout)

# Background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

# Subtitle / Category
sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11), Inches(0.4))
tf_sub = sub_box.text_frame
p_sub = tf_sub.paragraphs[0]
p_sub.text = "PHÂN TÍCH CHIẾN LƯỢC VẬN HÀNH ĐỐI TÁC"
p_sub.font.name = 'Arial'
p_sub.font.size = Pt(13)
p_sub.font.bold = True
p_sub.font.color.rgb = ORANGE

# Title
title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.7), Inches(11.3), Inches(1.5))
tf_title = title_box.text_frame
tf_title.word_wrap = True
p_title = tf_title.paragraphs[0]
p_title.text = "CƠ CHẾ XẾP HẠNG & PHÚC LỢI ĐẶC QUYỀN\nRANK × LAYER × AHABENEFITS V2.0"
p_title.font.name = 'Arial'
p_title.font.size = Pt(38)
p_title.font.bold = True
p_title.font.color.rgb = WHITE

# Brief / Description
desc_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(10), Inches(1.0))
tf_desc = desc_box.text_frame
tf_desc.word_wrap = True
p_desc = tf_desc.paragraphs[0]
p_desc.text = "Tái cấu trúc hệ thống xếp hạng tài xế dựa trên chất lượng DQS, phân tầng quyền hạn đăng ký ca theo vùng (Layer) và tối ưu hóa P&L bằng cơ chế AhaPoints."
p_desc.font.name = 'Arial'
p_desc.font.size = Pt(14)
p_desc.font.color.rgb = RGBColor(191, 219, 254) # Light blue-gray

# Footer info
footer_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(10), Inches(0.5))
tf_foot = footer_box.text_frame
p_foot = tf_foot.paragraphs[0]
p_foot.text = "AHAMOVE DRIVER MANAGEMENT TEAM | BÁO CÁO CHIẾN LƯỢC"
p_foot.font.name = 'Arial'
p_foot.font.size = Pt(10)
p_foot.font.bold = True
p_foot.font.color.rgb = WHITE

# ==========================================
# SLIDE 2: BỐI CẢNH & ĐIỂM NGHẼN (DESCRIPTIVE & DIAGNOSTIC)
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_header(slide, "Bối Cảnh Vận Hành & Điểm Nghẽn Cốt Lõi")

# Left Column: Hiện trạng & Bối cảnh
card_left_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.6), Inches(5.8), Inches(4.5))
card_left_bg.fill.solid()
card_left_bg.fill.fore_color.rgb = LIGHT_BG
card_left_bg.line.color.rgb = BORDER_COLOR

left_title = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.4), Inches(0.5))
left_title.text_frame.word_wrap = True
p = left_title.text_frame.paragraphs[0]
p.text = "🔍 Phân Tích Hiện Trạng (Descriptive)"
p.font.name = 'Arial'
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY

left_content = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.4), Inches(3.6))
left_content.text_frame.word_wrap = True
tf = left_content.text_frame
tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

bullets_left = [
    ("Hệ thống Ranking cũ rời rạc:", " Đánh giá tài xế dựa trên các chỉ số AR-FR-Rating độc lập, chưa phản ánh thực chất chất lượng phục vụ."),
    ("Thiếu sự phân hoá đặc quyền:", " Tài xế chạy sản lượng cao nhưng chất lượng kém vẫn được giữ rank cao, gây mất động lực phấn đấu."),
    ("Phân bổ ca thủ công:", " Cơ chế đăng ký ca chưa tạo được sự ưu tiên rõ nét cho nhóm tài xế nòng cốt có hiệu suất cao nhất."),
    ("Áp lực P&L nặng nề:", " Ngân sách hỗ trợ đảm bảo thu nhập (Promo/Incentives) chi trả bằng tiền mặt cố định, tốn kém nhưng hiệu quả điều phối kém.")
]

for title, desc in bullets_left:
    p = tf.add_paragraph()
    p.space_after = Pt(10)
    run1 = p.add_run()
    run1.text = "• " + title
    run1.font.bold = True
    run1.font.size = Pt(12)
    run1.font.color.rgb = DARK_TEXT
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(12)
    run2.font.color.rgb = GRAY_TEXT

# Right Column: Điểm nghẽn chẩn đoán
card_right_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.8), Inches(4.5))
card_right_bg.fill.solid()
card_right_bg.fill.fore_color.rgb = LIGHT_BG
card_right_bg.line.color.rgb = BORDER_COLOR

right_title = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.4), Inches(0.5))
right_title.text_frame.word_wrap = True
p = right_title.text_frame.paragraphs[0]
p.text = "⚡ Điểm Nghẽn Vận Hành (Diagnostic)"
p.font.name = 'Arial'
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ORANGE

right_content = slide.shapes.add_textbox(Inches(7.0), Inches(2.3), Inches(5.4), Inches(3.6))
right_content.text_frame.word_wrap = True
tf = right_content.text_frame
tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

bullets_right = [
    ("Khoảng trống chất lượng dịch vụ:", " Tài xế vi phạm quy chuẩn dịch vụ nhưng vẫn duy trì thăng hạng do thiếu thang điểm chất lượng tích hợp DQS."),
    ("Thiếu hụt cung giờ cao điểm (Peak-hour):", " Xảy ra tình trạng hụt tài xế cục bộ tại các Minizone (khu vực bán kính ngắn ≤ 4km) vào các giờ vàng."),
    ("Mất cân bằng Cung - Cầu:", " Không có cơ chế tự điều phối để tài xế tự nguyện lấp đầy các ca cần thiết, buộc phải tăng promo để ép tài xế chạy."),
    ("Rủi ro gian lận gia tăng:", " Quy trình tính thưởng ca thủ công trên Excel tạo kẽ hở gian lận điểm thưởng và sai sót số liệu vận hành.")
]

for title, desc in bullets_right:
    p = tf.add_paragraph()
    p.space_after = Pt(10)
    run1 = p.add_run()
    run1.text = "• " + title
    run1.font.bold = True
    run1.font.size = Pt(12)
    run1.font.color.rgb = DARK_TEXT
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(12)
    run2.font.color.rgb = GRAY_TEXT

# Insight footer bar
insight_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.3), Inches(12.0), Inches(0.7))
insight_bg.fill.solid()
insight_bg.fill.fore_color.rgb = RGBColor(239, 246, 255) # Light blue
insight_bg.line.color.rgb = RGBColor(191, 219, 254)

insight_text = slide.shapes.add_textbox(Inches(0.8), Inches(6.35), Inches(11.6), Inches(0.6))
insight_text.text_frame.word_wrap = True
p = insight_text.text_frame.paragraphs[0]
p.margin_left = p.margin_top = p.margin_right = p.margin_bottom = 0
run1 = p.add_run()
run1.text = "📌 Insight cốt lõi: "
run1.font.bold = True
run1.font.size = Pt(11.5)
run1.font.color.rgb = NAVY
run2 = p.add_run()
run2.text = "Để giải quyết triệt để, Ahamove cần tích hợp DQS vào lõi xét hạng, biến 'Ca hoạt động' thành tài nguyên khan hiếm được phân quyền đăng ký sớm theo Rank, và dùng điểm thưởng AhaPoints làm lực kéo tự nhiên."
run2.font.size = Pt(11.5)
run2.font.color.rgb = NAVY_LIGHT

# ==========================================
# SLIDE 3: MA TRẬN RANK × LAYER (PRESCRIPTIVE - PROPOSED MATRIX)
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_header(slide, "Ma Trận Quy Chuẩn: Rank × Layer × Quyền Lợi Ca")

# Subtitle of the slide
sub_tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.45), Inches(12.0), Inches(0.4))
sub_tb.text_frame.word_wrap = True
p_sub = sub_tb.text_frame.paragraphs[0]
p_sub.text = "Quy định điều kiện xét hạng (DQS), phân quyền Layer, khung giờ mở cổng đăng ký ca và hệ sinh thái phúc lợi đặc quyền."
p_sub.font.size = Pt(12)
p_sub.font.italic = True
p_sub.font.color.rgb = GRAY_TEXT

# Table
left = Inches(0.6)
top = Inches(1.9)
width = Inches(12.0)
height = Inches(4.5)

rows = 5
cols = 7
table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Set Column Widths
table.columns[0].width = Inches(1.8) # Rank
table.columns[1].width = Inches(2.2) # Layer ưu tiên
table.columns[2].width = Inches(1.5) # Khung giờ đăng ký
table.columns[3].width = Inches(1.3) # AhaPoints
table.columns[4].width = Inches(1.4) # Voucher xăng
table.columns[5].width = Inches(1.8) # Đội trưởng hỗ trợ
table.columns[6].width = Inches(2.0) # Đặc quyền khác

headers = ["Rank", "Layer Ưu Tiên", "Mở Đăng Ký Ca", "AhaPoints", "Xăng/EV", "Đội Trưởng", "Đặc Quyền Khác"]
for col_idx, text in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Arial'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE

# Table data
data = [
    [
        "💎 R1 Elite\n(Kim Cương | 15%)",
        "L2 Minizone (≤4km)\n(Cascade sang L3)",
        "00:00 - 10:00\n(Ngày 1)",
        "×1.5",
        "50k/tháng",
        "✅ Trưởng nhóm",
        "Bảo hiểm tai nạn,\nEV Priority Slot 1, Full-day"
    ],
    [
        "🥇 R2 Active\n(Hạng Vàng | 35%)",
        "L3 Mediumzone (≤8km)\n(Cascade sang L4)",
        "10:00 - 14:00\n(Ngày 1)",
        "×1.3",
        "30k/tháng",
        "✅ Trưởng nhóm",
        "Quà sinh nhật,\nEV Priority Slot 2, Full-day"
    ],
    [
        "🥈 R3 Standard\n(Hạng Bạc | 35%)",
        "L4 Bigzone\n(Cascade sang L5)",
        "Sau 14:00\n(Ngày 1)",
        "×1.1",
        "❌ Không",
        "✅ Phó nhóm",
        "Voucher data 5GB/tháng,\nHỗ trợ vá xe lưu động"
    ],
    [
        "👤 Unranked\n(Mới / Tự do | 15%)",
        "L6 MASS\n(Không cascade)",
        "Ngày 2+\n(Slot thừa)",
        "×1.0",
        "❌ Không",
        "❌ Không",
        "Tích điểm đổi thưởng cơ bản"
    ]
]

# Populate Table Data
for row_idx, row_data in enumerate(data):
    for col_idx, cell_value in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = cell_value
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Arial'
        p.font.size = Pt(10.5)
        
        # Format left-most cells differently
        if col_idx == 0:
            p.alignment = PP_ALIGN.LEFT
            p.font.bold = True
        elif col_idx == 1:
            p.alignment = PP_ALIGN.LEFT
            
        # Color formatting based on rank rows
        cell.fill.solid()
        if row_idx == 0:
            cell.fill.fore_color.rgb = ORANGE_BG
            p.font.color.rgb = DARK_TEXT
        elif row_idx == 1:
            cell.fill.fore_color.rgb = RGBColor(239, 246, 255) # Light blue
            p.font.color.rgb = DARK_TEXT
        elif row_idx == 2:
            cell.fill.fore_color.rgb = LIGHT_BG
            p.font.color.rgb = DARK_TEXT
        else:
            cell.fill.fore_color.rgb = WHITE
            p.font.color.rgb = GRAY_TEXT

# Footer note on table
note_tb = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.5))
p_note = note_tb.text_frame.paragraphs[0]
p_note.text = "⚠️ Điều kiện xét Rank định kỳ hàng tháng (DQS - Driver Quality Score): R1 Elite yêu cầu DQS ≥ 80 | R2 & R3 yêu cầu DQS ≥ 75. Năng suất (Productivity) là điều kiện đủ."
p_note.font.size = Pt(10)
p_note.font.bold = True
p_note.font.color.rgb = RED

# ==========================================
# SLIDE 4: POINTS ECONOMY & LỰC ĐẨY VẬN HÀNH (PREDICTIVE & ANALYSIS)
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_header(slide, "Cơ Chế AhaPoints & Kinh Tế Điểm (Points Economy)")

# Left Card: Công thức & Nguyên lý
left_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.6), Inches(5.8), Inches(4.5))
left_bg.fill.solid()
left_bg.fill.fore_color.rgb = LIGHT_BG
left_bg.line.color.rgb = BORDER_COLOR

left_tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.4), Inches(4.1))
tf = left_tb.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

# Title
p = tf.paragraphs[0]
p.text = "🪙 Cơ Chế Tích Điểm (Earn Points)"
p.font.name = 'Arial'
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY
p.space_after = Pt(12)

# Formula block
formula_p = tf.add_paragraph()
formula_p.text = " earned_pts = round( round(trip_GSV ÷ 1.000) × layer_multiplier ) "
formula_p.font.name = 'Courier New'
formula_p.font.size = Pt(11)
formula_p.font.bold = True
formula_p.font.color.rgb = DARK_TEXT
formula_p.space_after = Pt(8)

formula_desc_p = tf.add_paragraph()
formula_desc_p.text = "• Ý nghĩa công thức: Tích điểm trực tiếp theo doanh thu chuyến xe (5.000đ thu nhập thực tế = 1 điểm). Điểm số cơ bản được nhân hệ số dựa trên Layer hoạt động thực tế của đơn hàng, không phụ thuộc vào Rank của tài xế."
formula_desc_p.font.size = Pt(12)
formula_desc_p.font.color.rgb = GRAY_TEXT
formula_desc_p.space_after = Pt(10)

rules_bullets = [
    ("• Ngưỡng đổi quà tối thiểu:", " 5.000 pts (tương đương voucher 50.000đ cash). Điểm reset cuối mỗi quý."),
    ("• Đặc quyền xăng & Bảo hiểm:", " R1 Kim Cương được tặng voucher xăng 50k/tháng & đổi điểm lấy bảo hiểm tai nạn tự nguyện (gói 10k/285 pts, gói 30k/857 pts). R2 nhận voucher xăng 30k/tháng."),
    ("• Chế độ kỷ luật và ĐBCL:", " Trừ thẳng -50 pts/lần vi phạm ĐBCL (điểm tối thiểu sàn bằng 0, không âm).")
]

for title, desc in rules_bullets:
    p = tf.add_paragraph()
    p.space_after = Pt(8)
    run1 = p.add_run()
    run1.text = title
    run1.font.bold = True
    run1.font.size = Pt(12)
    run1.font.color.rgb = DARK_TEXT
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(12)
    run2.font.color.rgb = GRAY_TEXT

# Right Card: Dự phóng năng suất & Tác động điều phối
right_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.8), Inches(4.5))
right_bg.fill.solid()
right_bg.fill.fore_color.rgb = LIGHT_BG
right_bg.line.color.rgb = BORDER_COLOR

right_tb = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.4), Inches(4.1))
tf_right = right_tb.text_frame
tf_right.word_wrap = True
tf_right.margin_left = tf_right.margin_top = tf_right.margin_right = tf_right.margin_bottom = 0

p = tf_right.paragraphs[0]
p.text = "📈 Dự Phóng Năng Suất Tích Lũy / Tháng"
p.font.name = 'Arial'
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ORANGE
p.space_after = Pt(12)

projections = [
    ("💎 R1 Elite (Layer L2 - EPH ~70k/h):", " Tích luỹ 420 pts/ca 4h $\\rightarrow$ Đạt ~9.240 pts/tháng (chạy 22 ngày/tháng). Quyền lợi vượt trội, dư sức đổi quà."),
    ("🥇 R2 Active (Layer L3 - EPH ~65k/h):", " Tích luỹ 338 pts/ca 4h $\\rightarrow$ Đạt ~7.436 pts/tháng (chạy 22 ngày/tháng)."),
    ("🥈 R3 Standard (Layer L4 - EPH ~60k/h):", " Tích luỹ 264 pts/ca 4h $\\rightarrow$ Đạt ~5.808 pts/tháng (chạy 22 ngày/tháng)."),
    ("👤 Unranked (Layer L6 - EPH ~55k/h):", " Tích luỹ 220 pts/ca 4h $\\rightarrow$ Đạt tối đa ~4.840 pts/tháng. Không bao giờ đạt ngưỡng đổi quà tối thiểu (5.000 pts).")
]

for title, desc in projections:
    p = tf_right.add_paragraph()
    p.space_after = Pt(8)
    run1 = p.add_run()
    run1.text = title
    run1.font.bold = True
    run1.font.size = Pt(11.5)
    
    # Color coding of ranks
    if "R1" in title:
        run1.font.color.rgb = ORANGE
    elif "R2" in title:
        run1.font.color.rgb = NAVY_LIGHT
    elif "R3" in title:
        run1.font.color.rgb = DARK_TEXT
    else:
        run1.font.color.rgb = GRAY_TEXT
        
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(11.5)
    run2.font.color.rgb = GRAY_TEXT

# Highlighted Insight at the bottom of Right Card
p_ins = tf_right.add_paragraph()
p_ins.space_before = Pt(10)
run_ins = p_ins.add_run()
run_ins.text = "💡 Tác động điều phối: "
run_ins.font.bold = True
run_ins.font.size = Pt(11)
run_ins.font.color.rgb = GREEN
run_ins2 = p_ins.add_run()
run_ins2.text = "Cơ chế này tạo lực đẩy thăng hạng tự thân cực mạnh. Tài xế bắt buộc phải cải thiện chất lượng DQS và duy trì hoạt động tích cực để thăng hạng nhằm tiếp cận ca sớm và Layer có hệ số nhân cao (L2, L3)."
run_ins2.font.size = Pt(11)
run_ins2.font.italic = True
run_ins2.font.color.rgb = GRAY_TEXT

# Save Presentation
output_path = "/Users/ts-1148/Desktop/Pulu-workspace/Output/Ahamove/01. STRATEGY & PLANNING/driver-ranking/2026-07-driver-ranking-layer-benefits.pptx"
prs.save(output_path)
print(f"PowerPoint slide deck successfully generated and saved to: {output_path}")
